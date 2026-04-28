import io
import base64
import zipfile
from pathlib import Path


def convert_document(data: bytes, filename: str) -> dict:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return _convert_pdf(data)
    elif ext in (".docx", ".doc"):
        return _convert_docx(data)
    elif ext in (".pptx", ".ppt"):
        return _convert_pptx(data)
    else:
        return {"markdown": "", "images": [], "error": f"Unsupported format: {ext}"}


def _convert_pdf(data: bytes) -> dict:
    import fitz  # pymupdf

    doc = fitz.open(stream=data, filetype="pdf")
    md_parts = []
    images = []

    for page_num, page in enumerate(doc):
        blocks = page.get_text("blocks")  # (x0,y0,x1,y1,text,block_no,block_type)
        page_lines = []
        for b in blocks:
            if b[6] == 0:  # text block
                page_lines.append(b[4].strip())
        text = "\n\n".join(l for l in page_lines if l)
        if text:
            md_parts.append(f"## Page {page_num + 1}\n\n{text}")

        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            img_bytes = base_image["image"]
            ext = base_image["ext"]
            img_name = f"page{page_num + 1}_img{img_index + 1}.{ext}"
            images.append({
                "filename": img_name,
                "base64": base64.b64encode(img_bytes).decode(),
                "mime": f"image/{ext}",
            })

    return {"markdown": "\n\n".join(md_parts), "images": images}


def _convert_pptx(data: bytes) -> dict:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    md_parts = []
    images = []
    img_counter = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_lines = [f"## Slide {slide_num}"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if shape.is_placeholder and shape.placeholder_format:
                        ph_idx = shape.placeholder_format.idx
                        if ph_idx == 0:  # title
                            slide_lines[0] = f"## Slide {slide_num}: {text}"
                            continue
                    slide_lines.append(text)

            if shape.shape_type == 13:  # PICTURE
                img_counter += 1
                img_bytes = shape.image.blob
                ext = shape.image.ext
                img_name = f"slide{slide_num}_img{img_counter}.{ext}"
                images.append({
                    "filename": img_name,
                    "base64": base64.b64encode(img_bytes).decode(),
                    "mime": f"image/{ext}",
                })

        md_parts.append("\n\n".join(slide_lines))

    return {"markdown": "\n\n---\n\n".join(md_parts), "images": images}


def _convert_docx(data: bytes) -> dict:
    from docx import Document

    doc = Document(io.BytesIO(data))
    md_parts = []
    images = []
    img_counter = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            md_parts.append("")
            continue
        style = para.style.name if para.style else ""
        if style.startswith("Heading 1"):
            md_parts.append(f"# {text}")
        elif style.startswith("Heading 2"):
            md_parts.append(f"## {text}")
        elif style.startswith("Heading 3"):
            md_parts.append(f"### {text}")
        elif style.startswith("List"):
            md_parts.append(f"- {text}")
        else:
            md_parts.append(text)

    # Extract images from relationships
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            img_counter += 1
            img_data = rel.target_part.blob
            content_type = rel.target_part.content_type
            ext = content_type.split("/")[-1]
            img_name = f"image{img_counter}.{ext}"
            images.append({
                "filename": img_name,
                "base64": base64.b64encode(img_data).decode(),
                "mime": content_type,
            })

    return {"markdown": "\n\n".join(md_parts), "images": images}
